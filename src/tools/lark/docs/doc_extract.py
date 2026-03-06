"""文档文本提取模块

支持从以下文件格式中提取纯文本：
  - PDF  (.pdf)      → PyPDF2
  - PPTX (.pptx)     → python-pptx
  - DOCX (.docx)     → python-docx
  - XLSX (.xlsx)     → openpyxl
  - TXT/CSV/JSON等   → 直接解码
"""

import io
import os
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# 纯文本类扩展名（直接 decode）
_PLAIN_TEXT_EXTS = {".txt", ".csv", ".json", ".md", ".yaml", ".yml", ".xml", ".html", ".htm", ".log", ".srt"}


def extract_text(file_bytes: bytes, file_name: str) -> Optional[str]:
    """根据文件扩展名自动选择提取器，返回纯文本。

    Args:
        file_bytes: 文件二进制数据
        file_name: 文件名（用于判断格式）

    Returns:
        提取的文本内容，失败返回 None
    """
    ext = os.path.splitext(file_name)[1].lower()

    try:
        if ext == ".pdf":
            return _extract_pdf(file_bytes)
        elif ext == ".pptx":
            return _extract_pptx(file_bytes)
        elif ext == ".docx":
            return _extract_docx(file_bytes)
        elif ext == ".xlsx":
            return _extract_xlsx(file_bytes)
        elif ext in _PLAIN_TEXT_EXTS:
            return _extract_plain_text(file_bytes)
        else:
            logger.info("不支持的文件格式: %s (%s)", file_name, ext)
            return None
    except Exception as e:
        logger.error("提取文本失败: %s, error=%s", file_name, e, exc_info=True)
        return None


def _extract_pdf(file_bytes: bytes) -> str:
    """从 PDF 提取文本。"""
    from PyPDF2 import PdfReader

    reader = PdfReader(io.BytesIO(file_bytes))
    parts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            parts.append(text.strip())

    result = "\n\n".join(parts)
    logger.info("PDF 提取完成: %d 页, %d 字符", len(reader.pages), len(result))
    return result


def _extract_pptx(file_bytes: bytes) -> str:
    """从 PPTX 提取文本（按幻灯片分组）。"""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    parts = []

    for i, slide in enumerate(prs.slides, 1):
        slide_texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_texts.append(text)
            # 表格中的文本
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_texts:
                        slide_texts.append(" | ".join(row_texts))

        if slide_texts:
            parts.append(f"[Slide {i}]\n" + "\n".join(slide_texts))

    result = "\n\n".join(parts)
    logger.info("PPTX 提取完成: %d 页, %d 字符", len(prs.slides), len(result))
    return result


def _extract_docx(file_bytes: bytes) -> str:
    """从 DOCX 提取文本。"""
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)

    # 表格中的文本
    for table in doc.tables:
        for row in table.rows:
            row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_texts:
                parts.append(" | ".join(row_texts))

    result = "\n".join(parts)
    logger.info("DOCX 提取完成: %d 段落, %d 字符", len(doc.paragraphs), len(result))
    return result


def _extract_xlsx(file_bytes: bytes) -> str:
    """从 XLSX 提取文本（每个 sheet 的单元格内容）。"""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        sheet_lines = [f"[Sheet: {sheet_name}]"]
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            line = " | ".join(cells).strip()
            if line and line != "|":
                sheet_lines.append(line)
        if len(sheet_lines) > 1:
            parts.append("\n".join(sheet_lines))

    wb.close()
    result = "\n\n".join(parts)
    logger.info("XLSX 提取完成: %d sheets, %d 字符", len(wb.sheetnames), len(result))
    return result


def _extract_plain_text(file_bytes: bytes) -> str:
    """解码纯文本文件。"""
    for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            continue
    return file_bytes.decode("utf-8", errors="replace")


def get_supported_extensions() -> set[str]:
    """返回支持提取文本的所有扩展名。"""
    return {".pdf", ".pptx", ".docx", ".xlsx"} | _PLAIN_TEXT_EXTS
